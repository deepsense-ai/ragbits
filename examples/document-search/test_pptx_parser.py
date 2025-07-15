# This is a temporary script for development purposes and PR testing.
# It will be removed before merging.

from __future__ import annotations

import asyncio
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ragbits.core.sources.local import LocalFileSource
from ragbits.document_search.documents.document import Document, DocumentMeta, DocumentType
from ragbits.document_search.ingestion.parsers.pptx.parser import PptxDocumentParser


async def create_dummy_pptx(file_path: str):
    """Creates a dummy PPTX file for testing."""
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    if title and title.has_text_frame:
        title.text_frame.text = "Test Presentation"
    if subtitle and subtitle.has_text_frame:
        subtitle.text_frame.text = "A presentation for testing the PPTX parser."

    # Slide 2: Text, Shape, and Hyperlink
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    if title_shape and title_shape.has_text_frame:
        title_shape.text_frame.text = "This is a slide with text, a shape, and a hyperlink."

    body_shape = shapes.placeholders[1]
    if body_shape and body_shape.has_text_frame:
        tf = body_shape.text_frame
        tf.text = "This is a bullet point."

        p = tf.add_paragraph()
        p.text = "This is a line with a "
        r = p.add_run()
        r.text = "hyperlink"
        if r.hyperlink:
            r.hyperlink.address = "https://www.google.com"

    # Slide 3: Image
    img_slide_layout = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(img_slide_layout)
    img_path = "packages/ragbits-core/tests/assets/img/test.png"
    if os.path.exists(img_path):
        left = top = Inches(1)
        slide3.shapes.add_picture(img_path, left, top)

    # Slide 4: With speaker notes
    notes_slide_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(notes_slide_layout)
    if slide4.has_notes_slide:
        notes_slide = slide4.notes_slide
        if notes_slide.notes_text_frame:
            text_frame = notes_slide.notes_text_frame
            text_frame.text = "These are speaker notes for slide 4."

    prs.save(file_path)


async def main():
    """Main function to test the PPTX parser."""
    pptx_file = "test_pptx.pptx"
    await create_dummy_pptx(pptx_file)

    try:
        document_meta = DocumentMeta(
            document_type=DocumentType.PPTX,
            source=LocalFileSource(path=Path(pptx_file)),
        )
        document = Document.from_document_meta(document_meta, Path(pptx_file))

        parser = PptxDocumentParser()
        elements = await parser.parse(document)

        print(f"--- Extracted {len(elements)} elements ---")
        for element in elements:
            print(f"Type: {element.element_type}")
            print(f"Content: {element.text_representation}")
            print(f"Location: {element.location}")
            print("-" * 20)

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    asyncio.run(main()) 