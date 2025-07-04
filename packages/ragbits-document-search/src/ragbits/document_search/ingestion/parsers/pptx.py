from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ragbits.document_search.documents.document import Document, DocumentType
from ragbits.document_search.documents.element import Element, ElementLocation, ImageElement, TextElement
from ragbits.document_search.ingestion.parsers.base import DocumentParser


class PptxDocumentParser(DocumentParser):
    """Parser that extracts content from PPTX files using *python-pptx*.

    The parser retrieves text from all textual shapes, table cells and slide notes, as well as
    the binary bytes of pictures embedded in the presentation. Each piece of data is converted
    into a corresponding :class:`~ragbits.document_search.documents.element.TextElement` or
    :class:`~ragbits.document_search.documents.element.ImageElement`.
    """

    supported_document_types = {DocumentType.PPTX}

    async def parse(self, document: Document) -> list[Element]:
        """Parse the given PPTX document.

        Args:
            document: The document to parse.

        Returns:
            A list of extracted elements.
        """
        self.validate_document_type(document.metadata.document_type)
        presentation = Presentation(str(document.local_path))
        elements: list[Element] = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_location = ElementLocation(page_number=slide_idx)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    if text and text.strip():
                        elements.append(
                            TextElement(
                                document_meta=document.metadata,
                                location=slide_location,
                                content=text.strip(),
                            )
                        )
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text
                            if cell_text and cell_text.strip():
                                elements.append(
                                    TextElement(
                                        document_meta=document.metadata,
                                        location=slide_location,
                                        content=cell_text.strip(),
                                    )
                                )
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    description = getattr(shape, "alt_text", None) or None
                    elements.append(
                        ImageElement(
                            document_meta=document.metadata,
                            location=slide_location,
                            image_bytes=image_bytes,
                            description=description,
                        )
                    )

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    elements.append(
                        TextElement(
                            document_meta=document.metadata,
                            location=slide_location,
                            content=notes_text.strip(),
                        )
                    )

        return elements
