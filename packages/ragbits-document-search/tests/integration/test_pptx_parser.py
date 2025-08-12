from pathlib import Path

import pytest
from pptx import Presentation
from pptx.shapes.group import GroupShape

from ragbits.document_search.documents.document import DocumentMeta
from ragbits.document_search.documents.element import TextElement
from ragbits.document_search.ingestion.parsers.pptx import PptxDocumentParser


def _normalize_whitespace(value: str) -> str:
    """Normalize whitespace for robust substring checks."""
    return " ".join(value.split())


@pytest.mark.asyncio
async def test_pptx_parser_callbacks_integration() -> None:
    """Validate PPTX callbacks using a sample presentation asset."""
    pptx_path = Path(__file__).parent.parent / "assets" / "pptx" / "sample_presentation.pptx"
    assert pptx_path.exists()

    presentation = Presentation(pptx_path.as_posix())

    expected_notes: list[str] = []
    for slide in presentation.slides:
        notes_slide = getattr(slide, "notes_slide", None)
        notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
        text = getattr(notes_text_frame, "text", None)
        if not text:
            continue
        text = text.strip()
        if text:
            expected_notes.append(text)

    expected_links: list[str] = [
        f"Link: {getattr(getattr(getattr(shape, 'click_action', None), 'hyperlink', None), 'address', None)}"
        for slide in presentation.slides
        for shape in slide.shapes
        if not isinstance(shape, GroupShape)
        and getattr(getattr(getattr(shape, "click_action", None), "hyperlink", None), "address", None)
    ]

    cp = presentation.core_properties
    expected_metadata_lines: list[str] = []
    for key in ["author", "title", "subject", "keywords", "category", "created", "modified"]:
        value = getattr(cp, key, None)
        if value is None:
            continue
        value_str = str(value).strip()
        if value_str:
            expected_metadata_lines.append(f"{key}: {value_str}")

    document = await DocumentMeta.from_local_path(pptx_path).fetch()

    parser_with = PptxDocumentParser()
    elements_with = await parser_with.parse(document)
    text_with = "\n".join(e.content for e in elements_with if isinstance(e, TextElement))
    images_with = [e for e in elements_with if not isinstance(e, TextElement)]

    parser_without = PptxDocumentParser(pptx_callbacks=[])
    elements_without = await parser_without.parse(document)
    text_without = "\n".join(e.content for e in elements_without if isinstance(e, TextElement))
    images_without = [e for e in elements_without if not isinstance(e, TextElement)]

    assert len(images_with) == len(images_without)
    assert len(text_with) >= len(text_without)

    normalized_with = _normalize_whitespace(text_with)
    for note in expected_notes:
        assert _normalize_whitespace(note) in normalized_with

    for link_text in expected_links:
        assert link_text in text_with
        assert link_text not in text_without

    for meta_line in expected_metadata_lines:
        if meta_line in text_with:
            assert meta_line not in text_without
